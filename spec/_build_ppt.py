# -*- coding: utf-8 -*-
"""Author a curated STRIDE Coach BP pitch deck (~16 slides) from the vision doc."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = pathlib.Path(__file__).resolve().parent
OUT = HERE / "STRIDE_COACH_PITCH_DECK.pptx"

GREEN = RGBColor(0x16, 0xA3, 0x4A)
DGREEN = RGBColor(0x14, 0x53, 0x2D)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF0, 0xFD, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def fill_rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def content_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    fill_rect(s, 0, 0, SW, Inches(0.16), GREEN)          # top accent bar
    tf = box(s, Inches(0.55), Inches(0.42), Inches(12.3), Inches(0.9))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title; _set(r, 27, DARK, True)
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        _set(r2, 13, GREEN, False)
    return s


def bullets(slide, items, top=Inches(1.7), left=Inches(0.7),
            width=Inches(12.0), size=15, gap=6):
    tf = box(slide, left, top, width, Inches(5.3))
    first = True
    for it in items:
        lvl = 0; txt = it
        if isinstance(it, tuple):
            lvl, txt = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        bullet = "▸ " if lvl == 0 else "· "
        r = p.add_run(); r.text = bullet + txt
        _set(r, size if lvl == 0 else size - 1.5,
             DARK if lvl == 0 else GRAY, lvl == 0)
    return tf


def table_slide(title, headers, rows, subtitle=None, col_widths=None,
                fsize=12, header_fsize=12):
    s = content_slide(title, subtitle)
    nrow, ncol = len(rows) + 1, len(headers)
    left, top = Inches(0.6), Inches(1.85)
    width, height = Inches(12.1), Inches(0.5) * nrow
    gt = s.shapes.add_table(nrow, ncol, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = GREEN
        c.margin_left = Inches(0.08); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        tf = c.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = h; _set(r, header_fsize, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.margin_left = Inches(0.08); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            r = tf.paragraphs[0].add_run(); r.text = val
            _set(r, fsize, DARK, j == 0)
    return s


# ---------- 1. TITLE ----------
s = prs.slides.add_slide(BLANK)
fill_rect(s, 0, 0, SW, SH, DGREEN)
fill_rect(s, 0, Inches(4.55), SW, Inches(0.06), GREEN)
tf = box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.0))
r = tf.paragraphs[0].add_run(); r.text = "STRIDE Coach"; _set(r, 54, WHITE, True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "产品愿景 & 商业计划 (BP)"
_set(r, 22, RGBColor(0xBB, 0xF7, 0xD0), False)
tf2 = box(s, Inches(0.95), Inches(4.8), Inches(11.5), Inches(1.0))
r = tf2.paragraphs[0].add_run()
r.text = "面向中国市场 · 数字均为假设，需早期实测校准"
_set(r, 13, RGBColor(0x86, 0xEF, 0xAC), False)

# ---------- 2. 问题 ----------
s = content_slide("1. 问题：亿级跑者，同一组痛点")
bullets(s, [
    "凭感觉训练：不懂周期化 / 负荷管理，易过度训练受伤，或长期原地踏步",
    "数据用不起来：手表采集了海量心率/配速/负荷/HRV，却不知如何转成训练决策",
    "专业指导难获得：线下私教贵且受地域 / 通勤限制；线上训练营群发，跟不上个体每日状态",
    "需求是整体的：训练 + 营养 + 恢复 + 伤病 + 选赛 + 装备 —— 没有一个连贯入口",
    (1, "两类人群都成立：① 有成绩目标的备赛跑者；② 科学健康跑 / 减脂的大众跑者"),
], size=16)

# ---------- 3. 愿景 ----------
s = content_slide("2. 愿景：全程伴跑的 AI 教练")
bullets(s, [
    "AI + 个体数据，首次让「专业教练级」个性化指导可规模化、可负担、随时随地",
    "不只是训练计划工具，而是记得住你身体状态 / 训练历史 / 目标的全程伴跑教练",
    "覆盖跑者从入坑到 PB（或从久坐到可持续健康跑）的完整生命周期",
    (1, "对跑者：把专业训练规划 + 数据分析 + 装备/比赛决策，降到月费可负担"),
    (1, "对品牌/赛事：用真实个体数据语境精准推荐，比泛投广告 ROI 高"),
], size=16)

# ---------- 4. 市场规模 ----------
table_slide(
    "3. 市场规模（TAM / SAM / SOM）",
    ["层", "口径", "规模", "年化营收空间（假设）"],
    [
        ["TAM", "全体中国跑步人群 × 跑步消费", "1.63 亿跑步人口 / 产业池 >4000 亿元", "双引擎理论可触达 4000 亿+ 消费池"],
        ["SAM", "有表 + 训练意图 + 付费力", "~1500–3000 万人（跑步人群 10–18%）", "纯订阅天花板 ~90–180 亿元/年"],
        ["SOM", "3 年现实可获取付费用户", "~15–90 万付费（SAM 渗透 1–3%）", "订阅 ARR ~0.7–4.2 亿元/年 + 导购"],
    ],
    subtitle="来源：中国田协《2024 路跑工作报告》等；人群交集 / 渗透率 / ARPU 为假设",
    col_widths=[1.0, 3.0, 4.0, 4.1], fsize=12.5,
)

# ---------- 5. 护城河 ----------
s = content_slide("4. 护城河：四层飞轮（不是数据，也不是公式）")
bullets(s, [
    "① 多设备数据归一化：工程基线，对标 RQ（非独有）—— 增量在覆盖华为/小米国产表",
    "② 大规模真实数据校准的算法精度：公式公开，准确度靠数据量 + 工程打磨，越用越难追平",
    "③ 周期化执行闭环：规划→执行→跟踪→调整，越用越贴合 —— Sigma/RQ 都还浅的一段",
    "④ 中文训练学 / 运动医学知识库：让伤病/营养/恢复敢答且可信，难以速成",
    (1, "四者咬合成数据飞轮：能跟它们咬合的能力扩展，价值都被放大（如个体语境导购）"),
], size=15)

# ---------- 6. 竞争格局 ----------
s = content_slide("5. 竞争格局：直接 vs 间接", "分水岭 = 是否已本地化争夺中文付费跑者")
bullets(s, [
    "直接竞品（同市场、同人群）：Sigma（超超世世·上海）、RQ 跑力（RunningQuotient）",
    (1, "目前唯一与我们争夺中国跑者的两家 —— 但训练系统的周期化对话闭环都还浅"),
    "品类标杆 / 间接（形态参考，不进中国）：Runna（全球最强）、TrainAsONE、Humango、Coopah",
    (1, "验证了 AI 教练模式成立，但无中文 / 无本土设备/支付/赛事，不与我们正面争夺"),
], size=15)

# ---------- 7. 竞品深拆 ----------
table_slide(
    "5. 竞品深拆：他们的短板就是我们的切口",
    ["竞品", "定位", "短板", "STRIDE 切口"],
    [
        ["Sigma", "记录审美先行的中文 AI 跑步", "周期化计划仍邀请制、知识库偏通识", "训练系统为核心 + 闭环深度"],
        ["RQ 跑力", "华语圈数据分析最严谨", "止步分析层、计划静态写死、无执行追踪", "对话入口 + 逐日自适应闭环"],
        ["Runna", "全球品类标杆（不进中国）", "只改配速不吃生理信号、无对话、伤病口碑差", "生理自适应 + 双向伤病感知教练"],
    ],
    col_widths=[1.5, 3.2, 4.2, 3.2], fsize=12,
)

# ---------- 8. 产品 ----------
s = content_slide("6. 产品：对话即入口")
bullets(s, [
    "一个对话框说任意诉求 → 系统识别意图并路由（建计划/调计划/状态/挪课/伤病/装备/选赛/订酒店）",
    "一条滚动会话贯穿连续旅程，记得住当前操作对象（目标赛事 / 计划 / 某一周）",
    "被动响应 + 主动教练 push（提醒/预警/督促）—— 主动 push 是留存与差异化核心",
    "App 优先 + Web（资深用户深度分析）",
    (1, "5 不变量：写操作先提案后确认 · 安全敏感不自动改 · 知识库可信 · 外部不编造 · 导购透明"),
], size=14.5)

# ---------- 9. 能力域 ----------
s = content_slide("7. 产品边界：11 能力域 × 3 层战略")
bullets(s, [
    "🟢 核心命脉（订阅引擎 R1）：0 计划 · 1 洞察 · 2 健康 · 3 营养 · 4 恢复 · 6 比赛策略 · 8 改课 · 10 知识",
    (1, "服务跑者全部训练需求，完整性本身就是留存护城河 —— 营养/恢复/改课不是加分项"),
    "🟡 外部世界（导购引擎 R2/R3/R4）：5 装备 · 6 比赛发现 · 7 酒店",
    (1, "接外部数据，用个体语境精准推荐；营养/恢复的实物导购也走这里"),
    "⚪ 社区（引流）：9 社交 —— 拉新与留存，后续品牌/广告",
], size=14.5)

# ---------- 10. 商业模式 ----------
s = content_slide("8. 商业模式：双引擎")
bullets(s, [
    "引擎一 · 订阅（recurring）：变现「完整训练服务」—— 覆盖越全，留存越牢，月费越值",
    "引擎二 · 导购分成（transactional）：变现「外部世界」—— 跑鞋/装备/补给 + 报名 + 酒店，按成交抽成",
    "两引擎共享同一资产：个体数据 —— 让订阅更专业、让导购更精准（数据越多两边都更值钱）",
    (1, "互为飞轮：订阅数据让导购更准；导购触点把免费用户拉向订阅"),
    (1, "导购对免费用户也变现 —— 即使不订阅也能产生收入"),
], size=15)

# ---------- 11. 收入流 ----------
table_slide(
    "8. 收入流 R1–R7",
    ["#", "收入流", "类型", "说明（假设）"],
    [
        ["R1", "个人订阅", "月/年费", "¥49/月，含跑步+力量个性化指导（主轴，估值锚）"],
        ["R2", "装备导购分成", "CPS 佣金", "跑鞋/装备联盟 5–15%，高客单竞速鞋（高毛利）"],
        ["R3", "比赛报名导流", "按单 / %", "固定 ¥/单 或报名费 %"],
        ["R4", "酒店/旅行分成", "OTA", "4–10% booking commission"],
        ["R5", "品牌合作 / 测评", "赞助费", "新品测评、官方训练伙伴"],
        ["R6", "真人教练 marketplace", "平台抽成", "撮合远程在线真人教练，抽 15–30%"],
        ["R7", "匿名聚合数据洞察", "B2B", "换鞋周期/装备偏好 → 品牌选品（隐私门控）"],
    ],
    col_widths=[0.7, 3.0, 2.0, 6.4], fsize=11.5,
)

# ---------- 12. 定价 ----------
table_slide(
    "9. 用户分层与定价（假设，待验证）",
    ["档", "价格", "目标人群"],
    [
        ["Free（引流）", "¥0 —— 数据同步、基础状态、有限对话、社区、被导购", "拉新 / 新手"],
        ["Pro（个人 AI 教练）", "¥49/月；6 月付 9 折、年付 8 折（约 ¥470）", "进阶 / 备赛跑者（主力付费）"],
        ["Elite / 备赛包", "¥99/月（或赛季包）+ 真人教练 + 深度报告", "严肃备赛 / 追 PB"],
    ],
    col_widths=[2.6, 6.5, 3.0], fsize=12.5,
)

# ---------- 13. 单位经济 ----------
s = table_slide(
    "10. 单位经济（base case，全负载毛利 30%）",
    ["指标", "假设值", "指标", "假设值"],
    [
        ["混合 ARPU", "¥45/月", "合计 LTV", "~¥300（订阅+导购）"],
        ["月流失 churn", "6% → 生命周期 ~16 月", "混合 CAC", "~¥130（2/3 自然流量）"],
        ["毛利率", "30%（含分摊研发）", "LTV/CAC", "~2.3（接近健康线）"],
        ["订阅毛利 LTV", "~¥216", "回本周期", "~10 月（健康线内）"],
    ],
    subtitle="三杠杆推向 LTV/CAC≥3：规模摊薄研发提毛利 · 压 churn · 维持自然流量+导购净增量",
    col_widths=[2.6, 3.5, 2.4, 3.6], fsize=12.5,
)

# ---------- 14. GTM ----------
s = content_slide("11. 获客 / GTM")
bullets(s, [
    "破冷启动：onboarding 一键同步主流手表历史数据 → Day-1 兑现价值（能力评估 + 首周计划）",
    "渠道：手表用户导入（最低门槛）· 小红书/B站/抖音 KOL · 赛事/跑团合作 · 社区裂变 · 应用商店/SEO",
    "免费增值漏斗：Free 拉新 → 高意图时刻（建赛季计划/深度诊断）转 Pro",
    (1, "增长飞轮：更多用户 → 更多数据 → 算法/推荐更准 → 留存更高 + 口碑扩散 → 更多用户"),
], size=15)

# ---------- 15. 路线 ----------
s = content_slide("12. 迭代路线（建议非承诺）")
bullets(s, [
    "P0 统一对话入口：意图识别路由，打通建/改计划 + 状态问答 ｜ 上线 R1 订阅墙",
    "P1 训练核心补全：挪课/跳课/反馈/改目标重排 ｜ 转化漏斗调优",
    "P2 身体 & 安全：伤病门控 + 重返跑步 + 力量诊断 ｜ Elite + R6 真人教练试点",
    "P3 数据加持咨询：营养/恢复/天气改课 ｜ R2 补给/恢复品导购",
    "P4 外部世界：比赛/装备/后勤 ｜ R2 跑鞋 + R3 报名 + R4 酒店 + R5 品牌",
    "P5 社区 & 知识 & 数据产品 ｜ R7 B2B 数据",
], size=14)

# ---------- 16. 开放问题 ----------
s = content_slide("13. 开放问题 & 下一步")
bullets(s, [
    "待验证：Pro 定价与付费意愿 · 导购转化率 · 单位经济三杠杆能否兑现",
    "待选型：赛事日历 / 跑鞋库 / 酒店外部数据接入",
    "待把关：伤病医疗边界（建议 vs 转诊）· 导购合规与信任 · R7 数据隐私",
    "BP 待补章节：团队 / 里程碑时间线 / 融资诉求 · 风险矩阵与缓解",
], size=15)

prs.save(str(OUT))
print(f"OK -> {OUT}  ({OUT.stat().st_size//1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
